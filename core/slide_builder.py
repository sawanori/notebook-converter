"""
PowerPoint slide generation module.

Creates PPTX files with background images and editable text boxes
based on OCR results.
"""
from __future__ import annotations

import io
import logging
import threading
from pathlib import Path
from typing import TYPE_CHECKING, List, Optional

from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from .geometry import calculate_slide_dimensions, px_to_emu

if TYPE_CHECKING:
    from .ocr_processor import TextBlock

logger = logging.getLogger(__name__)

# Font configuration
DEFAULT_FONT_NAME = "Yu Gothic UI"  # Japanese-compatible font
FALLBACK_FONT_NAME = "Arial"
DEFAULT_FONT_SIZE = Pt(11)


class SlideBuilder:
    """
    Builds PowerPoint presentations from images and text blocks.
    
    Creates slides with background images and positions text
    in editable text boxes based on OCR coordinates.
    """

    def __init__(self, dpi: int = 300):
        """
        Initialize the slide builder.

        Args:
            dpi: DPI of source images (default: 300).
        """
        self.dpi = dpi
        self.prs: Optional[Presentation] = None
        self._slide_width: int = 0
        self._slide_height: int = 0
        logger.info(f"SlideBuilder initialized with DPI: {dpi}")

    def create_presentation(self, first_image: Image.Image) -> None:
        """
        Initialize a new presentation.

        Sets slide dimensions based on the first image's aspect ratio.

        Args:
            first_image: First slide image to determine dimensions.
        """
        self.prs = Presentation()

        # Calculate slide dimensions from image
        width_emu, height_emu = calculate_slide_dimensions(
            first_image.width,
            first_image.height,
            self.dpi
        )

        self.prs.slide_width = width_emu
        self.prs.slide_height = height_emu
        self._slide_width = width_emu
        self._slide_height = height_emu

        logger.info(
            f"Presentation created: {first_image.width}x{first_image.height}px "
            f"-> {width_emu}x{height_emu} EMU"
        )

    def add_slide(
        self,
        image: Image.Image,
        text_blocks: List[TextBlock],
        stop_event: Optional[threading.Event] = None
    ) -> None:
        """
        Add a slide with background image and text boxes.

        Args:
            image: Background image for the slide.
            text_blocks: List of text blocks to add as text boxes.
            stop_event: Optional event to signal cancellation.

        Raises:
            RuntimeError: If presentation not initialized.
        """
        if stop_event and stop_event.is_set():
            logger.debug("Slide addition cancelled")
            return

        if self.prs is None:
            raise RuntimeError(
                "Presentation not initialized. Call create_presentation() first."
            )

        # Get blank layout (index 6), with fallback
        slide_layout = self._get_blank_layout()
        slide = self.prs.slides.add_slide(slide_layout)

        # Add background image
        self._add_background_image(slide, image)

        # Add text boxes
        for block in text_blocks:
            if stop_event and stop_event.is_set():
                logger.debug("Text box addition cancelled")
                return
            self._add_text_box(slide, block)

        logger.debug(f"Added slide with {len(text_blocks)} text boxes")

    def _get_blank_layout(self):
        """
        Get a blank slide layout.

        Tries index 6 first (standard blank), falls back to last layout.

        Returns:
            Slide layout object.
        """
        try:
            # Standard blank layout is usually at index 6
            if len(self.prs.slide_layouts) > 6:
                return self.prs.slide_layouts[6]
            else:
                # Fallback to last layout
                return self.prs.slide_layouts[-1]
        except IndexError:
            # Ultimate fallback: first layout
            return self.prs.slide_layouts[0]

    def _add_background_image(
        self,
        slide,
        image: Image.Image
    ) -> None:
        """
        Add background image covering the entire slide.

        Args:
            slide: Slide object to add image to.
            image: PIL Image to use as background.
        """
        # Convert image to bytes
        image_stream = io.BytesIO()
        
        # Use PNG for lossless quality
        image.save(image_stream, format="PNG")
        image_stream.seek(0)

        # Add picture at position (0, 0) covering full slide
        slide.shapes.add_picture(
            image_stream,
            Emu(0),
            Emu(0),
            self._slide_width,
            self._slide_height
        )

    def _add_text_box(
        self,
        slide,
        block: TextBlock
    ) -> None:
        """
        Add an editable text box to the slide.

        Positions the text box based on OCR coordinates and
        makes it transparent so background shows through.

        Args:
            slide: Slide object to add text box to.
            block: TextBlock with text and position information.
        """
        # Convert coordinates to EMU
        left = px_to_emu(block.left, self.dpi)
        top = px_to_emu(block.top, self.dpi)
        width = px_to_emu(block.width, self.dpi)
        height = px_to_emu(block.height, self.dpi)

        # Add text box shape
        textbox = slide.shapes.add_textbox(
            Emu(left),
            Emu(top),
            Emu(width),
            Emu(height)
        )

        # Configure text frame
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.auto_size = None  # Don't auto-resize

        # Add text to first paragraph
        p = tf.paragraphs[0]
        p.text = block.text
        p.alignment = PP_ALIGN.LEFT

        # Set font properties
        for run in p.runs:
            run.font.name = DEFAULT_FONT_NAME
            run.font.size = DEFAULT_FONT_SIZE

        # If no runs (empty paragraph with text set directly)
        if not p.runs and p.text:
            p.font.name = DEFAULT_FONT_NAME
            p.font.size = DEFAULT_FONT_SIZE

        # Make text box transparent
        textbox.fill.background()  # No fill
        textbox.line.fill.background()  # No border

    def save(self, output_path: Path) -> None:
        """
        Save the presentation to a file.

        Args:
            output_path: Path to save the PPTX file.

        Raises:
            RuntimeError: If no presentation to save.
            IOError: If file cannot be written.
        """
        if self.prs is None:
            raise RuntimeError("No presentation to save.")

        try:
            self.prs.save(str(output_path))
            logger.info(f"Presentation saved to: {output_path}")
        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            raise IOError(f"Failed to save presentation: {e}") from e

    @property
    def slide_count(self) -> int:
        """Get the number of slides in the presentation."""
        if self.prs is None:
            return 0
        return len(self.prs.slides)
